import docx
import pypandoc
import openpyxl
from pdfminer.high_level import extract_text
from markdownify import markdownify as md
import ebooklib
from ebooklib import epub
from pptx import Presentation

def convert_docx_to_md(docx_file):
    """Convert .docx file to .md"""
    doc = docx.Document(docx_file)
    md_content = ""

    for para in doc.paragraphs:
        md_content += para.text + "\n\n"

    return md_content

def convert_xlsx_to_md(xlsx_file):
    """Convert .xlsx file to .md (Excel)"""
    wb = openpyxl.load_workbook(xlsx_file)
    md_content = ""

    for sheet in wb.sheetnames:
        md_content += f"# {sheet}\n"
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            md_content += "\t".join(str(cell) for cell in row) + "\n"
        md_content += "\n\n"

    return md_content

def convert_pptx_to_md(pptx_file):
    """Convert .pptx file to .md (PowerPoint)"""
    prs = Presentation(pptx_file)
    md_content = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                md_content += shape.text + "\n"
        md_content += "\n\n"

    return md_content

def convert_pdf_to_md(pdf_file):
    """Convert .pdf file to .md"""
    text = extract_text(pdf_file)
    return text

def convert_html_to_md(html_file):
    """Convert .html file to .md"""
    with open(html_file, "r", encoding="utf-8") as f:
        html_content = f.read()
    return md(html_content)

def convert_epub_to_md(epub_file):
    """Convert .epub file to .md"""
    book = epub.read_epub(epub_file)
    md_content = ""
    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            md_content += item.get_body_content().decode("utf-8") + "\n\n"
    return md_content

def convert_mobi_to_md(mobi_file):
    """Convert .mobi file to .md using pypandoc"""
    return pypandoc.convert_file(mobi_file, 'md')

def convert_odt_to_md(odt_file):
    """Convert .odt file to .md using pypandoc"""
    return pypandoc.convert_file(odt_file, 'md')

def convert_rtf_to_md(rtf_file):
    """Convert .rtf file to .md using pypandoc"""
    return pypandoc.convert_file(rtf_file, 'md')

def convert_latex_to_md(latex_file):
    """Convert .latex file to .md using pypandoc"""
    return pypandoc.convert_file(latex_file, 'md')

def convert_txt_to_md(txt_file):
    """Convert .txt file to .md"""
    with open(txt_file, "r", encoding="utf-8") as f:
        return f.read()

def convert_to_md(input_file):
    """Main function to detect file type and convert to Markdown"""
    file_extension = input_file.lower().split('.')[-1]

    if file_extension == 'docx':
        return convert_docx_to_md(input_file)
    elif file_extension == 'pdf':
        return convert_pdf_to_md(input_file)
    elif file_extension == 'html':
        return convert_html_to_md(input_file)
    elif file_extension == 'epub':
        return convert_epub_to_md(input_file)
    elif file_extension == 'mobi':
        return convert_mobi_to_md(input_file)
    elif file_extension == 'odt':
        return convert_odt_to_md(input_file)
    elif file_extension == 'rtf':
        return convert_rtf_to_md(input_file)
    elif file_extension == 'latex' or file_extension == 'tex':
        return convert_latex_to_md(input_file)
    elif file_extension == 'txt':
        return convert_txt_to_md(input_file)
    elif file_extension == 'xlsx':
        return convert_xlsx_to_md(input_file)
    elif file_extension == 'pptx':
        return convert_pptx_to_md(input_file)
    else:
        return convert_other_formats(input_file)

def convert_other_formats(input_file):
    """Convert formats not explicitly handled by other functions"""
    try:
        output = pypandoc.convert_file(input_file, 'md')
        return output
    except Exception as e:
        raise ValueError(f"Unsupported file format or conversion error: {e}")

def save_md(content, output_file):
    """Save the converted content to a .md file"""
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write(content)